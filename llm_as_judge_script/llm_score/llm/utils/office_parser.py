"""
Office Document Parser
Supports Word, Excel, PowerPoint, PDF document formats
"""
import os
import tempfile
import subprocess
from io import BytesIO
from pathlib import Path
from typing import Dict, Any, List, Optional

# Office file parsing dependencies
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import pandas as pd
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

try:
    import PyPDF2
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    import openpyxl
    OPENPYXL_SUPPORT = True
except ImportError:
    OPENPYXL_SUPPORT = False

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_SUPPORT = True
except ImportError:
    PDF2IMAGE_SUPPORT = False


class OfficeParser:
    """Office Document Parser"""
    
    SUPPORTED_FORMATS = {
        '.docx': 'Word Document',
        '.doc': 'Word Document (Legacy)',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation (Legacy)',
        '.xlsx': 'Excel Workbook',
        '.xls': 'Excel Workbook (Legacy)',
        '.pdf': 'PDF Document'
    }
    
    @staticmethod
    def _estimate_tokens(text: str) -> int:
        """
        估算文本的token数量（约4个字符=1个token）
        
        Args:
            text: 输入文本
            
        Returns:
            估算的token数量
        """
        return len(text) // 4 + 1
    
    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        """Check if file is a supported Office format"""
        return Path(file_path).suffix.lower() in cls.SUPPORTED_FORMATS
    
    @classmethod
    def get_supported_formats(cls) -> Dict[str, str]:
        """Get supported file formats"""
        return cls.SUPPORTED_FORMATS.copy()
    
    @classmethod
    async def parse_document(cls, file_data: bytes, file_path: str) -> str:
        """
        Parse Office document
        
        Args:
            file_data: Document file data
            file_path: File path
            
        Returns:
            Parsed document content
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == '.docx':
            return await cls._parse_docx(file_data, file_path)
        elif file_ext == '.doc':
            return await cls._parse_doc(file_data, file_path)
        elif file_ext == '.pptx':
            return await cls._parse_pptx(file_data, file_path)
        elif file_ext == '.ppt':
            return await cls._parse_ppt(file_data, file_path)
        elif file_ext == '.xlsx':
            return await cls._parse_xlsx(file_data, file_path)
        elif file_ext == '.xls':
            return await cls._parse_xls(file_data, file_path)
        elif file_ext == '.pdf':
            return await cls._parse_pdf(file_data, file_path)
        else:
            return f"Unsupported Office document format: {file_ext}"
    
    @classmethod
    async def _parse_docx(cls, file_data: bytes, file_path: str) -> str:
        """Parse DOCX file"""
        if not DOCX_SUPPORT:
            return "python-docx library required to parse DOCX files"
        
        try:
            doc = Document(BytesIO(file_data))
            content = []
            
            # Extract paragraph content
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Check if it's a heading
                    if paragraph.style.name.startswith('Heading'):
                        level = paragraph.style.name.replace('Heading ', '')
                        content.append(f"Heading {level}: {text}")
                    else:
                        content.append(text)
            
            # Extract table content
            for i, table in enumerate(doc.tables):
                content.append(f"\nTable {i+1}:")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cell for cell in cells):  # Ensure not an empty row
                        content.append("  " + " | ".join(cells))
            
            result = "\n".join(content)
            return result if result.strip() else "Document content is empty"
            
        except Exception as e:
            return f"Failed to parse DOCX file: {str(e)}"
    
    @classmethod
    async def _parse_doc(cls, file_data: bytes, file_path: str) -> str:
        """Parse DOC file (legacy Word format)"""
        return "DOC file parsing requires additional library support, recommend converting to DOCX format"
    
    @classmethod
    async def _parse_pptx(cls, file_data: bytes, file_path: str) -> str:
        """Parse PPTX file"""
        if not PPTX_SUPPORT:
            return "python-pptx library required to parse PPTX files"
        
        try:
            prs = Presentation(BytesIO(file_data))
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract slide title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_text.append(f"Title: {title_text}")
                
                # Extract other text content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # Avoid duplicating title
                        if not (slide.shapes.title and text == slide.shapes.title.text.strip()):
                            slide_text.append(text)
                
                # Extract table content
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        slide_text.append("Table:")
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            if any(cell for cell in cells):
                                slide_text.append("  " + " | ".join(cells))
                
                if slide_text:
                    slides_content.append(f"Slide {i}:\n" + "\n".join(slide_text))
                else:
                    slides_content.append(f"Slide {i}: (no text content)")
            
            result = "\n\n".join(slides_content)
            return result if result.strip() else "Presentation content is empty"
            
        except Exception as e:
            return f"Failed to parse PPTX file: {str(e)}"
    
    @classmethod
    async def _parse_ppt(cls, file_data: bytes, file_path: str) -> str:
        """Parse PPT file (legacy PowerPoint format)"""
        return "PPT file parsing requires additional library support, recommend converting to PPTX format"
    
    @classmethod
    async def _parse_pptx_as_images(cls, file_data: bytes, file_path: str, dpi: int = 150, max_slides: int = 50) -> List[bytes]:
        """
        Render PPTX file as list of images
        
        Args:
            file_data: PPTX file data
            file_path: File path (for getting filename)
            dpi: Image rendering DPI
            max_slides: Maximum number of slides to render
            
        Returns:
            List of image bytes, each element is a PNG image of a slide
        """
        if not PDF2IMAGE_SUPPORT:
            raise Exception("pdf2image library required to render PPTX as images")
        
        temp_pptx = None
        temp_pdf = None
        temp_dir = None
        
        try:
            # Create temporary directory
            temp_dir = tempfile.mkdtemp()
            
            # Save PPTX to temporary file
            temp_pptx = os.path.join(temp_dir, "presentation.pptx")
            with open(temp_pptx, 'wb') as f:
                f.write(file_data)
            
            # Use LibreOffice to convert PPTX to PDF
            temp_pdf = os.path.join(temp_dir, "presentation.pdf")
            
            # Try different libreoffice commands
            libreoffice_commands = [
                'libreoffice',
                'soffice',
                '/usr/bin/libreoffice',
                '/usr/bin/soffice'
            ]
            
            conversion_success = False
            for cmd in libreoffice_commands:
                try:
                    result = subprocess.run(
                        [cmd, '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, temp_pptx],
                        capture_output=True,
                        timeout=60,
                        check=False
                    )
                    if result.returncode == 0 and os.path.exists(temp_pdf):
                        conversion_success = True
                        break
                except (FileNotFoundError, subprocess.TimeoutExpired):
                    continue
            
            if not conversion_success:
                raise Exception(
                    "Unable to convert PPTX to PDF using LibreOffice. Please ensure LibreOffice is installed: "
                    "sudo apt-get install libreoffice (Linux) or brew install libreoffice (macOS)"
                )
            
            # Read PDF file
            with open(temp_pdf, 'rb') as f:
                pdf_data = f.read()
            
            # Use pdf2image to convert PDF to images
            images = convert_from_bytes(pdf_data, dpi=dpi, fmt='png')
            
            # Limit number of slides
            if len(images) > max_slides:
                images = images[:max_slides]
            
            # Convert PIL Images to bytes
            image_bytes_list = []
            for img in images:
                img_byte_arr = BytesIO()
                img.save(img_byte_arr, format='PNG')
                image_bytes_list.append(img_byte_arr.getvalue())
            
            return image_bytes_list
            
        except Exception as e:
            raise Exception(f"Failed to render PPTX as images: {str(e)}")
        
        finally:
            # Clean up temporary files
            try:
                if temp_pptx and os.path.exists(temp_pptx):
                    os.remove(temp_pptx)
                if temp_pdf and os.path.exists(temp_pdf):
                    os.remove(temp_pdf)
                if temp_dir and os.path.exists(temp_dir):
                    os.rmdir(temp_dir)
            except Exception:
                pass  # Ignore cleanup errors
    
    @classmethod
    async def _parse_xlsx(cls, file_data: bytes, file_path: str, max_tokens: int = 8000) -> str:
        """
        Parse XLSX file with dynamic row truncation based on token limit
        
        Args:
            file_data: Excel file data
            file_path: File path
            max_tokens: Maximum tokens allowed (default: 8000)
            
        Returns:
            Parsed Excel content as string
        """
        if not EXCEL_SUPPORT:
            return "pandas and openpyxl libraries required to parse XLSX files"
        
        try:
            # Read all worksheets
            excel_file = pd.ExcelFile(BytesIO(file_data))
            content = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    # Filter out completely empty rows and columns
                    df = df.dropna(how='all').dropna(axis=1, how='all')
                    
                    if not df.empty:
                        worksheet_header = f"Worksheet: {sheet_name}"
                        
                        # Keep all columns - generate header with all column names
                        column_names = list(df.columns)
                        header_str = " | ".join(str(col) for col in column_names)
                        header_line = f"Columns: {header_str}"
                        
                        # Calculate tokens for worksheet header and column header
                        current_tokens = cls._estimate_tokens(worksheet_header) + cls._estimate_tokens(header_line)
                        sheet_content = [worksheet_header, header_line]
                        
                        # Add rows dynamically based on token limit
                        rows_added = 0
                        total_rows = len(df)
                        
                        for idx, row in df.iterrows():
                            # Convert row to string format
                            row_values = [str(val) if pd.notna(val) else '' for val in row]
                            row_str = " | ".join(row_values)
                            
                            # Estimate tokens for this row
                            row_tokens = cls._estimate_tokens(row_str)
                            
                            # Check if adding this row would exceed token limit
                            if current_tokens + row_tokens > max_tokens:
                                break
                            
                            sheet_content.append(row_str)
                            current_tokens += row_tokens
                            rows_added += 1
                        
                        # Add truncation notice if needed
                        if rows_added < total_rows:
                            truncation_notice = f"(Showing first {rows_added} rows out of {total_rows} total rows due to token limit)"
                            sheet_content.append(truncation_notice)
                        
                        content.extend(sheet_content)
                        content.append("")
                    else:
                        content.append(f"Worksheet: {sheet_name} (empty)")
                        content.append("")
                        
                except Exception as e:
                    content.append(f"Worksheet: {sheet_name} - Parsing failed: {str(e)}")
                    content.append("")
            
            result = "\n".join(content)
            return result if result.strip() else "Excel file content is empty"
            
        except Exception as e:
            return f"Failed to parse XLSX file: {str(e)}"
    
    @classmethod
    async def _parse_xls(cls, file_data: bytes, file_path: str, max_tokens: int = 8000) -> str:
        """
        Parse XLS file (legacy Excel format) with dynamic row truncation based on token limit
        
        Args:
            file_data: Excel file data
            file_path: File path
            max_tokens: Maximum tokens allowed (default: 8000)
            
        Returns:
            Parsed Excel content as string
        """
        if not EXCEL_SUPPORT:
            return "pandas and xlrd libraries required to parse XLS files"
        
        try:
            # Use xlrd engine to read legacy Excel files
            excel_file = pd.ExcelFile(BytesIO(file_data), engine='xlrd')
            content = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name, engine='xlrd')
                    
                    # Filter out completely empty rows and columns
                    df = df.dropna(how='all').dropna(axis=1, how='all')
                    
                    if not df.empty:
                        worksheet_header = f"Worksheet: {sheet_name}"
                        
                        # Keep all columns - generate header with all column names
                        column_names = list(df.columns)
                        header_str = " | ".join(str(col) for col in column_names)
                        header_line = f"Columns: {header_str}"
                        
                        # Calculate tokens for worksheet header and column header
                        current_tokens = cls._estimate_tokens(worksheet_header) + cls._estimate_tokens(header_line)
                        sheet_content = [worksheet_header, header_line]
                        
                        # Add rows dynamically based on token limit
                        rows_added = 0
                        total_rows = len(df)
                        
                        for idx, row in df.iterrows():
                            # Convert row to string format
                            row_values = [str(val) if pd.notna(val) else '' for val in row]
                            row_str = " | ".join(row_values)
                            
                            # Estimate tokens for this row
                            row_tokens = cls._estimate_tokens(row_str)
                            
                            # Check if adding this row would exceed token limit
                            if current_tokens + row_tokens > max_tokens:
                                break
                            
                            sheet_content.append(row_str)
                            current_tokens += row_tokens
                            rows_added += 1
                        
                        # Add truncation notice if needed
                        if rows_added < total_rows:
                            truncation_notice = f"(Showing first {rows_added} rows out of {total_rows} total rows due to token limit)"
                            sheet_content.append(truncation_notice)
                        
                        content.extend(sheet_content)
                        content.append("")
                    else:
                        content.append(f"Worksheet: {sheet_name} (empty)")
                        content.append("")
                        
                except Exception as e:
                    content.append(f"Worksheet: {sheet_name} - Parsing failed: {str(e)}")
                    content.append("")
            
            result = "\n".join(content)
            return result if result.strip() else "Excel file content is empty"
            
        except Exception as e:
            return f"Failed to parse XLS file: {str(e)}"
    
    @classmethod
    async def _parse_pdf(cls, file_data: bytes, file_path: str) -> str:
        """Parse PDF file"""
        if not PDF_SUPPORT:
            return "PyPDF2 library required to parse PDF files"
        
        try:
            pdf_reader = PyPDF2.PdfReader(BytesIO(file_data))
            content = []
            
            # Get PDF basic information
            num_pages = len(pdf_reader.pages)
            content.append(f"PDF document info: {num_pages} pages total")
            
            # Extract document metadata
            if pdf_reader.metadata:
                metadata = pdf_reader.metadata
                if metadata.get('/Title'):
                    content.append(f"Title: {metadata['/Title']}")
                if metadata.get('/Author'):
                    content.append(f"Author: {metadata['/Author']}")
                if metadata.get('/Subject'):
                    content.append(f"Subject: {metadata['/Subject']}")
            
            content.append("\nDocument content:")
            
            # Extract text content from each page
            max_pages = 20  # Limit pages to process
            for i, page in enumerate(pdf_reader.pages[:max_pages]):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        content.append(f"\n--- Page {i+1} ---")
                        # Clean text format
                        cleaned_text = ' '.join(page_text.split())
                        content.append(cleaned_text)
                except Exception as e:
                    content.append(f"\n--- Page {i+1} (extraction failed: {str(e)}) ---")
            
            if num_pages > max_pages:
                content.append(f"\n... ({num_pages - max_pages} more pages not shown)")
            
            result = "\n".join(content)
            return result if result.strip() else "PDF file unable to extract text content"
            
        except Exception as e:
            return f"Failed to parse PDF file: {str(e)}"
    
    @classmethod
    async def get_document_info(cls, file_data: bytes, file_path: str) -> Dict[str, Any]:
        """
        Get document basic information
        
        Args:
            file_data: Document file data
            file_path: File path
            
        Returns:
            Document information dictionary
        """
        file_ext = Path(file_path).suffix.lower()
        info = {
            'file_name': Path(file_path).name,
            'file_type': cls.SUPPORTED_FORMATS.get(file_ext, 'Unknown format'),
            'file_size': len(file_data)
        }
        
        try:
            if file_ext == '.docx' and DOCX_SUPPORT:
                doc = Document(BytesIO(file_data))
                info.update({
                    'paragraphs_count': len(doc.paragraphs),
                    'tables_count': len(doc.tables)
                })
            
            elif file_ext == '.pptx' and PPTX_SUPPORT:
                prs = Presentation(BytesIO(file_data))
                info.update({
                    'slides_count': len(prs.slides)
                })
            
            elif file_ext == '.xlsx' and EXCEL_SUPPORT:
                excel_file = pd.ExcelFile(BytesIO(file_data))
                info.update({
                    'sheets_count': len(excel_file.sheet_names),
                    'sheet_names': excel_file.sheet_names
                })
            
            elif file_ext == '.pdf' and PDF_SUPPORT:
                pdf_reader = PyPDF2.PdfReader(BytesIO(file_data))
                info.update({
                    'pages_count': len(pdf_reader.pages)
                })
                
                if pdf_reader.metadata:
                    metadata = pdf_reader.metadata
                    info.update({
                        'title': metadata.get('/Title', ''),
                        'author': metadata.get('/Author', ''),
                        'subject': metadata.get('/Subject', '')
                    })
        
        except Exception as e:
            info['error'] = f"Failed to get document info: {str(e)}"
        
        return info
    
    @classmethod
    async def parse_pptx_for_images(cls, file_data: bytes, file_path: str, dpi: int = 150, max_slides: int = 50) -> List[bytes]:
        """
        Parse PPTX file as list of images (for multimodal models)
        
        Args:
            file_data: PPTX file data
            file_path: File path
            dpi: Image rendering DPI (default 150)
            max_slides: Maximum number of slides to render (default 50)
            
        Returns:
            List of image bytes, each element is a PNG image of a slide
        """
        return await cls._parse_pptx_as_images(file_data, file_path, dpi, max_slides)
